from __future__ import annotations

import io
import logging
from typing import Any, List, Optional, Tuple

from ...parser.models import (
    ParsedCell,
    ParsedTable,
    Provenance,
    TableRef,
    TextChunk,
    TextStyle,
)

logger = logging.getLogger(__name__)


class PptxExtractor:
    """
    Extract text and tables from PowerPoint files using python-pptx.

    When *image_extractor* is provided (ENABLE_VISION=true), picture shapes
    on each slide are described via Gemini Vision and their descriptions are
    included in the slide text so visual content (diagrams, screenshots,
    charts) is not lost.
    """

    def __init__(self, image_extractor: Optional[Any] = None) -> None:
        self._image_extractor = image_extractor

    def extract(
        self,
        data: bytes,
        filename: str,
        provenance: Provenance,
    ) -> Tuple[str, List[ParsedTable]]:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(io.BytesIO(data))
        text_parts: List[str] = []
        tables: List[ParsedTable] = []
        table_index = 0

        for slide_num, slide in enumerate(prs.slides):
            slide_texts: List[str] = []

            for shape in slide.shapes:
                # ── Text frames ───────────────────────────────────────────
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            slide_texts.append(line)

                # ── Embedded images ───────────────────────────────────────
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and self._image_extractor:
                    description = self._describe_shape_image(
                        shape, slide_num, filename, provenance
                    )
                    if description:
                        slide_texts.append(
                            f"[Image on slide {slide_num + 1}]: {description}"
                        )

                # ── Tables ────────────────────────────────────────────────
                if shape.has_table:
                    tbl = shape.table
                    chain = provenance.table_chain.copy()
                    chain.append(TableRef(table_index=table_index, row=None, col=None))
                    table_prov = Provenance(
                        page_id=provenance.page_id,
                        page_title=provenance.page_title,
                        page_date=provenance.page_date,
                        table_chain=chain,
                        attachment_id=provenance.attachment_id,
                        attachment_type=provenance.attachment_type,
                    )

                    headers: List[str] = []
                    cells_matrix: List[List[ParsedCell]] = []

                    for row_idx, row in enumerate(tbl.rows):
                        row_cells: List[ParsedCell] = []
                        for col_idx, cell in enumerate(row.cells):
                            text = cell.text.strip()
                            cell_prov = Provenance(
                                page_id=provenance.page_id,
                                page_title=provenance.page_title,
                                page_date=provenance.page_date,
                                table_chain=chain.copy(),
                                row=row_idx,
                                col=col_idx,
                                attachment_id=provenance.attachment_id,
                                attachment_type=provenance.attachment_type,
                            )
                            chunks = (
                                [TextChunk(content=text, style=TextStyle.NORMAL, provenance=cell_prov)]
                                if text else []
                            )
                            row_cells.append(ParsedCell(
                                row=row_idx, col=col_idx,
                                text_chunks=chunks,
                                provenance=cell_prov,
                            ))
                            if text:
                                slide_texts.append(text)

                        if row_idx == 0:
                            headers = [
                                c.text_chunks[0].content if c.text_chunks else ""
                                for c in row_cells
                            ]
                        cells_matrix.append(row_cells)

                    tables.append(ParsedTable(
                        table_index=table_index,
                        headers=headers,
                        cells=cells_matrix,
                        provenance=table_prov,
                        raw_html=f"[PPTX slide {slide_num + 1} table]",
                    ))
                    table_index += 1

            if slide_texts:
                text_parts.append(f"[Slide {slide_num + 1}]")
                text_parts.extend(slide_texts)

        return "\n".join(text_parts), tables

    # ------------------------------------------------------------------
    # Vision helper
    # ------------------------------------------------------------------

    def _describe_shape_image(
        self,
        shape,
        slide_num: int,
        filename: str,
        provenance: Provenance,
    ) -> str:
        """Extract image bytes from a picture shape and describe via vision."""
        try:
            img_bytes = shape.image.blob
            if not img_bytes:
                return ""
            description, _ = self._image_extractor.extract(
                data=img_bytes,
                filename=filename,
                provenance=provenance,
            )
            return description
        except Exception as exc:
            logger.warning(
                "Vision failed for image on slide %d of '%s': %s",
                slide_num + 1, filename, exc,
            )
            return ""
